"""
モノシリ ファイルテキスト抽出モジュール
各ファイル形式からテキストとページ/スライドのメタデータを抽出する。
外部送信なし・完全ローカル処理。
"""
from __future__ import annotations
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# 対応ファイル拡張子
SUPPORTED_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".doc",
    ".pptx",
    ".xlsx", ".xls",
    ".txt", ".csv", ".md",
    ".dwg", ".dxf",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp",
}

# テキストを完全抽出できる形式
EXTRACTABLE_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".csv", ".md",
}

# ファイル名のみ記録する形式（テキスト抽出不可）
FILENAME_ONLY_EXTENSIONS: set[str] = {
    ".doc", ".xls",
    ".dwg", ".dxf",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp",
}

# 形式の表示名マッピング
EXTENSION_LABELS: dict[str, str] = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".doc": "Word（旧形式）",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".xls": "Excel（旧形式）",
    ".txt": "テキスト",
    ".csv": "CSV",
    ".md": "Markdown",
    ".dwg": "CAD（DWG）",
    ".dxf": "CAD（DXF）",
    ".jpg": "画像（JPG）",
    ".jpeg": "画像（JPEG）",
    ".png": "画像（PNG）",
    ".gif": "画像（GIF）",
    ".bmp": "画像（BMP）",
}


def extract_text(file_path: Path) -> tuple[list[dict], str | None]:
    """
    ファイルからテキストとメタデータを抽出する。

    Returns:
        (page_chunks, skip_reason)
        page_chunks: [{"text": str, "page": int|None, "slide": int|None}]
        skip_reason: スキップ理由（成功時はNone）
    """
    ext = file_path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        return [], f"未対応形式: {ext}"

    # ファイル名のみ記録する形式
    if ext in FILENAME_ONLY_EXTENSIONS:
        label = EXTENSION_LABELS.get(ext, ext)
        return [{"text": f"[{label}] {file_path.name}", "page": None, "slide": None}], None

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext == ".docx":
            return _extract_docx(file_path)
        elif ext == ".pptx":
            return _extract_pptx(file_path)
        elif ext == ".xlsx":
            return _extract_xlsx(file_path)
        elif ext in {".txt", ".csv", ".md"}:
            return _extract_text_file(file_path)
    except MemoryError:
        return [], "メモリ不足（ファイルが大きすぎます）"
    except Exception as e:
        logger.error(f"抽出エラー {file_path.name}: {e}")
        return [], f"抽出エラー: {str(e)[:120]}"

    return [], "未対応形式"


def _extract_pdf(file_path: Path) -> tuple[list[dict], None]:
    """PyMuPDFを使いページごとにテキスト抽出

    BUG-027修正: get_text("text") → get_text("blocks") + 座標ソート
    財務PDFの表構造（年度×指標）を保持するため、テキストブロックをY座標→X座標順に
    ソートし直す。これにより列単位での誤抽出を防ぎ、行単位の自然読み順を保持する。
    round(b[1] / 10) でY座標を10px単位に丸め、同一行の左右ブロックをX座標で並べる。
    """
    import fitz  # PyMuPDF

    page_chunks = []
    doc = fitz.open(str(file_path))
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            # blocks = [(x0, y0, x1, y1, text, block_no, block_type), ...]
            blocks = page.get_text("blocks")
            # テキストブロック(block_type==0)のみ、Y→X座標順でソート
            text_blocks = sorted(
                [b for b in blocks if b[6] == 0],
                key=lambda b: (round(b[1] / 10), b[0])  # 行単位(10px)でY→X
            )
            text = "\n".join(b[4].strip() for b in text_blocks if b[4].strip())
            if text:
                page_chunks.append({
                    "text": text,
                    "page": page_num + 1,
                    "slide": None,
                })
    finally:
        doc.close()

    return page_chunks, None


def _extract_docx(file_path: Path) -> tuple[list[dict], None]:
    """python-docxを使い全文テキスト抽出

    CHUNK-002修正: セクションヘッダーを各テキスト行にプレフィックスとして付与。
    Heading スタイルを検出して current_heading を更新し、後続の段落・テーブル行に
    「[見出し]」プレフィックスを付ける。チャンク分割後も「どのセクションか」が保持される。
    テーブルは行単位で「セル1 | セル2 | ...」形式に変換し、行の意味関係を保持する。
    結合セルの重複はXML要素IDで検出して排除する。
    """
    from docx import Document

    doc = Document(str(file_path))
    texts: list[str] = []
    current_heading: str = ""

    # 段落テキスト（見出し検出付き）
    for para in doc.paragraphs:
        stripped = para.text.strip()
        if not stripped:
            continue

        style_name = para.style.name if para.style else ""
        is_heading = (
            style_name.startswith("Heading")
            or "見出し" in style_name
            or style_name.lower().startswith("heading")
        )

        if is_heading:
            current_heading = stripped
            texts.append(stripped)
        else:
            if current_heading:
                texts.append(f"[{current_heading}] {stripped}")
            else:
                texts.append(stripped)

    # テーブル（行単位で構造化・見出しコンテキスト付与）
    for table in doc.tables:
        header_cells: list[str] = []
        for row_idx, row in enumerate(table.rows):
            # 結合セルの重複をXML要素IDで排除
            seen: set[int] = set()
            row_cells: list[str] = []
            for cell in row.cells:
                cell_id = id(cell._tc)
                if cell_id not in seen:
                    seen.add(cell_id)
                    row_cells.append(cell.text.strip())

            non_empty = [c for c in row_cells if c]
            if not non_empty:
                continue

            if row_idx == 0:
                # 最初の行をヘッダー候補として記録
                header_cells = row_cells
                row_line = " | ".join(non_empty)
            else:
                # ヘッダーと値を対応付けて「ヘッダー: 値」形式で出力
                if header_cells:
                    parts: list[str] = []
                    for header, value in zip(header_cells, row_cells):
                        if value:
                            parts.append(f"{header}: {value}" if header else value)
                    row_line = " | ".join(parts) if parts else " | ".join(non_empty)
                else:
                    row_line = " | ".join(non_empty)

            if current_heading:
                texts.append(f"[{current_heading}] {row_line}")
            else:
                texts.append(row_line)

    full_text = "\n".join(texts)
    if not full_text:
        return [], None

    return [{"text": full_text, "page": None, "slide": None}], None


def _extract_pptx(file_path: Path) -> tuple[list[dict], None]:
    """python-pptxを使いスライドごとにテキスト抽出"""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    page_chunks = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            # テキストフレームを持つシェイプ
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            elif hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            # テーブル
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            texts.append(t)

        if texts:
            page_chunks.append({
                "text": "\n".join(texts),
                "page": None,
                "slide": slide_num,
            })

    return page_chunks, None


def _extract_xlsx(file_path: Path) -> tuple[list[dict], None]:
    """openpyxlを使い行列関係を保持してテキスト抽出

    CHUNK-003修正: 最初の非空行をヘッダーとして読み取り、以降のデータ行を
    「ヘッダー: 値 | ヘッダー: 値 ...」形式で出力する。
    シート名をプレフィックスとして付与することで、チャンク後も列の意味が保持される。
    例: [損益計算書] 売上高: 7,623,374 | 営業利益: 523,812
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
    all_texts: list[str] = []

    for sheet in wb.worksheets:
        sheet_name: str = sheet.title or ""
        prefix: str = f"[{sheet_name}] " if sheet_name else ""
        header_row: list[str] = []
        sheet_texts: list[str] = []

        for row in sheet.iter_rows(values_only=True):
            # None → 空文字列に統一
            cell_values = [
                str(c).strip() if (c is not None and str(c).strip().lower() != "none") else ""
                for c in row
            ]

            # 全空行はスキップ
            if not any(cell_values):
                continue

            if not header_row:
                # 最初の非空行をヘッダー候補に
                header_row = cell_values
                header_line = " | ".join(v for v in header_row if v)
                if header_line:
                    sheet_texts.append(f"{prefix}{header_line}")
            else:
                # データ行: ヘッダーと対応付けて「ヘッダー: 値」形式で出力
                parts: list[str] = []
                for i, value in enumerate(cell_values):
                    if not value:
                        continue
                    if i < len(header_row) and header_row[i]:
                        parts.append(f"{header_row[i]}: {value}")
                    else:
                        parts.append(value)
                if parts:
                    sheet_texts.append(f"{prefix}" + " | ".join(parts))

        all_texts.extend(sheet_texts)

    wb.close()

    if not all_texts:
        return [], None

    return [{"text": "\n".join(all_texts), "page": None, "slide": None}], None


def _extract_text_file(file_path: Path) -> tuple[list[dict], str | None]:
    """テキスト・CSV・Markdownファイルを読み込む"""
    # 文字コードを順番に試みる
    for encoding in ("utf-8", "utf-8-sig", "shift_jis", "cp932", "euc_jp", "latin-1"):
        try:
            text = file_path.read_text(encoding=encoding).strip()
            if text:
                return [{"text": text, "page": None, "slide": None}], None
            return [], None  # 空ファイル
        except UnicodeDecodeError:
            continue
        except Exception as e:
            return [], f"読み込みエラー: {str(e)[:80]}"

    return [], "文字コードを判定できませんでした"


def scan_folder(folder_path: Path) -> list[Path]:
    """
    フォルダ内の対象ファイルをすべて列挙する（サブフォルダ含む）
    隠しファイル・隠しフォルダはスキップする。
    """
    files: list[Path] = []
    try:
        for f in folder_path.rglob("*"):
            # 隠しファイル・フォルダをスキップ
            if any(part.startswith(".") for part in f.parts):
                continue
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS:
                files.append(f)
    except PermissionError as e:
        logger.warning(f"アクセス権限なし: {folder_path} — {e}")

    return sorted(files)


def estimate_index_time(file_count: int) -> str:
    """ファイル数から処理時間を見積もる（目安）"""
    if file_count == 0:
        return "0秒"
    # 経験則: 1ファイル約2秒（Embedding含む）
    seconds = file_count * 2
    if seconds < 60:
        return f"約{seconds}秒"
    elif seconds < 3600:
        return f"約{seconds // 60}分"
    else:
        return f"約{seconds // 3600}時間{(seconds % 3600) // 60}分"
